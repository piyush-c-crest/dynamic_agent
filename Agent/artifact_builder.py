"""
artifact_builder.py
===================
Document and presentation builders for the Dynamic Agent Orchestration system.

Extracted from dynamic_langgraph_backend.py to keep format-conversion code
separate from the LangGraph orchestration logic.

All functions here are pure transformations: they take text/JSON in and
produce a python-docx Document / python-pptx Presentation / bytes out.
They have no dependency on LangGraph, LLMs, or the tool registry.

Supported output formats
------------------------
- Markdown -> .docx  (_md_to_docx)     requires: python-docx
- JSON slides -> .pptx (_build_pptx)   requires: python-pptx
- Markdown -> .pdf   (via weasyprint, invoked in DynamicToolRegistry)
- raw SVG  -> .svg   (plain write, no helper needed)

The _ARTIFACT_KINDS set is the authoritative list consumed by the
create_artifact tool in DynamicToolRegistry to validate the `kind` argument.

Public API
----------
_ARTIFACT_KINDS                 -- {"doc", "pdf", "slides", "image"}
_md_to_docx(content)            -- convert Markdown str -> python-docx Document
_build_pptx(slides)             -- convert list[dict] -> python-pptx Presentation
_add_inline_runs(paragraph, text) -- internal helper, exposed for testing
"""

import re

# Optional heavy dependencies -- imported lazily so the module is importable
# even if the packages aren't installed. The create_artifact tool checks for
# None and returns a helpful error message instead of crashing.
try:
    from docx import Document as _DocxDocument
except ImportError:
    _DocxDocument = None

try:
    from pptx import Presentation as _PptxPresentation
except ImportError:
    _PptxPresentation = None


# ---------------------------------------------------------------------------
# Artifact kind registry
# ---------------------------------------------------------------------------

# All valid values for the `kind` parameter of the create_artifact tool.
_ARTIFACT_KINDS = {"doc", "pdf", "slides", "image"}


# ---------------------------------------------------------------------------
# Inline Markdown parser (bold / italic / code -> docx runs)
# ---------------------------------------------------------------------------

_INLINE_MD_RE = re.compile(
    r"\*\*(?P<bold1>.+?)\*\*"
    r"|__(?P<bold2>.+?)__"
    r"|\*(?P<italic1>.+?)\*"
    r"|_(?P<italic2>.+?)_"
    r"|`(?P<code>.+?)`"
)


def _add_inline_runs(paragraph, text: str) -> None:
    """Splits a line on inline markdown (bold/italic/code) into styled runs
    so a downloaded .docx shows real bold/italic instead of literal '**x**'."""
    pos = 0
    for match in _INLINE_MD_RE.finditer(text):
        if match.start() > pos:
            paragraph.add_run(text[pos:match.start()])
        if match.group("bold1") is not None:
            paragraph.add_run(match.group("bold1")).bold = True
        elif match.group("bold2") is not None:
            paragraph.add_run(match.group("bold2")).bold = True
        elif match.group("italic1") is not None:
            paragraph.add_run(match.group("italic1")).italic = True
        elif match.group("italic2") is not None:
            paragraph.add_run(match.group("italic2")).italic = True
        elif match.group("code") is not None:
            paragraph.add_run(match.group("code")).font.name = "Courier New"
        pos = match.end()
    if pos < len(text) or pos == 0:
        paragraph.add_run(text[pos:])


# ---------------------------------------------------------------------------
# Markdown table helpers
# ---------------------------------------------------------------------------

def _parse_md_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _is_md_table_separator(line: str) -> bool:
    cells = _parse_md_table_row(line)
    return bool(cells) and all(re.fullmatch(r":?-+:?", c) for c in cells)


def _add_docx_table(document, header: list[str], rows: list[list[str]]) -> None:
    table = document.add_table(rows=1 + len(rows), cols=len(header))
    table.style = "Table Grid"
    for col, text in enumerate(header):
        paragraph = table.rows[0].cells[col].paragraphs[0]
        _add_inline_runs(paragraph, text)
        for run in paragraph.runs:
            run.bold = True
    for r, row in enumerate(rows, start=1):
        for col, text in enumerate(row):
            if col < len(header):
                _add_inline_runs(table.rows[r].cells[col].paragraphs[0], text)


# ---------------------------------------------------------------------------
# Markdown -> .docx converter
# ---------------------------------------------------------------------------

def _md_to_docx(content: str):
    """Small, deliberately non-exhaustive Markdown -> docx converter: enough
    structure for agent-authored reports (headings, bullets, paragraphs,
    tables, inline bold/italic/code), not a full CommonMark implementation.

    Raises RuntimeError if python-docx is not installed.
    """
    if _DocxDocument is None:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    document = _DocxDocument()
    lines = content.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped:
            i += 1
            continue
        # Markdown table (detected by a separator row immediately after the header)
        if stripped.startswith("|") and i + 1 < len(lines) and _is_md_table_separator(lines[i + 1]):
            header = _parse_md_table_row(stripped)
            i += 2
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(_parse_md_table_row(lines[i]))
                i += 1
            _add_docx_table(document, header, rows)
            continue
        # ATX-style headings
        if stripped.startswith("### "):
            _add_inline_runs(document.add_heading("", level=3), stripped[4:])
        elif stripped.startswith("## "):
            _add_inline_runs(document.add_heading("", level=2), stripped[3:])
        elif stripped.startswith("# "):
            _add_inline_runs(document.add_heading("", level=1), stripped[2:])
        # Bullet points
        elif stripped.startswith(("- ", "* ")):
            _add_inline_runs(document.add_paragraph(style="List Bullet"), stripped[2:])
        # Fallback: normal paragraph
        else:
            _add_inline_runs(document.add_paragraph(), stripped)
        i += 1
    return document


# ---------------------------------------------------------------------------
# JSON slides -> .pptx builder
# ---------------------------------------------------------------------------

def _build_pptx(slides: list[dict]):
    """Convert a list of slide dicts (each with 'title' and optional 'bullets')
    into a python-pptx Presentation using the 'Title and Content' layout.

    Raises RuntimeError if python-pptx is not installed.
    """
    if _PptxPresentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    presentation = _PptxPresentation()
    layout = presentation.slide_layouts[1]  # "Title and Content"
    for entry in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = entry["title"]
        bullets = entry.get("bullets") or []
        if bullets:
            body = slide.placeholders[1].text_frame
            body.text = bullets[0]
            for bullet in bullets[1:]:
                body.add_paragraph().text = bullet
    return presentation
