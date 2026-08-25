from langgraph.graph import StateGraph, START, END
from typing import TypedDict, Annotated, Any, Callable
from langchain_core.messages import (
    BaseMessage, HumanMessage, AIMessage, SystemMessage, ToolMessage,
)
from langgraph.checkpoint.sqlite import SqliteSaver
from langgraph.graph.message import add_messages
from langchain_community.utilities import GoogleSerperAPIWrapper
from langchain_core.tools import tool, Tool, BaseTool
from dotenv import load_dotenv
import os
import sys
import sqlite3
import requests
import json
import re
import base64
import uuid
import operator
import contextvars
import subprocess
import venv
from pathlib import Path
from datetime import datetime
from langchain_openai import ChatOpenAI
from langchain_core.tools import StructuredTool
from pydantic import create_model, Field
from langsmith import traceable

import document_pipeline as docpipe

# Optional dependencies for the create_artifact tool. Each is only needed for
# the artifact kinds that use it (docx for "doc", pptx for "slides",
# weasyprint for "pdf") -- missing one doesn't break the other kinds, it
# just makes that specific kind return a clear "pip install X" error instead
# of failing to import at module load time.
try:
    import markdown as _markdown
except ImportError:
    _markdown = None
try:
    from docx import Document as _DocxDocument
except ImportError:
    _DocxDocument = None
try:
    from pptx import Presentation as _PptxPresentation
except ImportError:
    _PptxPresentation = None
try:
    import weasyprint as _weasyprint
except ImportError:
    _weasyprint = None

load_dotenv()

# ---------- terminal output -> file logging ----------
class _Tee:
    """Mirrors every write to both the original stream and a log file.

    Wrapping sys.stdout/sys.stderr this way captures ALL console output —
    including every existing print() call throughout this module — without
    needing to touch each call site individually.
    """
    def __init__(self, *streams):
        self.streams = streams

    def write(self, data):
        for s in self.streams:
            try:
                s.write(data)
                s.flush()
            except Exception:
                pass  # never let a logging failure break the app

    def flush(self):
        for s in self.streams:
            try:
                s.flush()
            except Exception:
                pass

    def isatty(self):
        return any(getattr(s, "isatty", lambda: False)() for s in self.streams)


def _setup_console_logging(log_dir: str = "logs") -> str:
    """Redirect stdout/stderr so everything printed to the terminal is also written to a file."""
    os.makedirs(log_dir, exist_ok=True)
    log_path = os.path.join(log_dir, f"run_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")
    log_file = open(log_path, "a", encoding="utf-8")  # utf-8 to safely hold any generated content
    sys.stdout = _Tee(sys.stdout, log_file)
    sys.stderr = _Tee(sys.stderr, log_file)
    return log_path


LOG_FILE_PATH = _setup_console_logging()
print(f"🗒️ Logging console output to: {LOG_FILE_PATH}")


def _log(tag: str, msg: str = "") -> None:
    """Single-line, timestamped log entry. `tag` identifies which part of the
    workflow this line is about (WORKFLOW / AI-REQUEST / AI-REPLY / TOOL-CALL /
    TOOL-RESULT / FINAL-RESPONSE / ...) so the console/log file can be scanned
    or grepped for just one stage of a run."""
    ts = datetime.now().strftime("%H:%M:%S.%f")[:-3]
    print(f"[{ts}] [{tag}] {msg}")


def _log_block(tag: str, title: str, body: str, max_chars: int = 2000) -> None:
    """Same as _log but for multi-line payloads (prompts, AI replies, tool
    output) -- prints a header line followed by the (possibly truncated)
    body so long content doesn't spam the single-line log stream."""
    body = body if body is not None else ""
    truncated = len(body) > max_chars
    shown = body[:max_chars] + (f"\n... [truncated, {len(body) - max_chars} more chars]" if truncated else "")
    _log(tag, f"{title}:")
    for line in shown.splitlines() or [""]:
        print(f"    {line}")

os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
os.environ["LANGCHAIN_PROJECT"] = "Dynamic-LangGraph-Backend"
# Requires LANGCHAIN_API_KEY in .env

if not os.environ.get("LANGCHAIN_API_KEY"):
    print("⚠️ LANGCHAIN_API_KEY not set — LangSmith tracing will not appear in the dashboard.")

MAX_RETRIES = 2
MAX_TASKS = 6
AUTO_TOOL_LIMIT = 2  # max new tools an agent's tool-selection step may auto-create per call
MAX_TOOL_VALIDATION_RETRIES = 3  # repair attempts if a generated tool fails its smoke test
MAX_TOOL_CALLS_PER_TASK = 3  # max consecutive tool calls per task to prevent infinite loops
MAX_READ_CHARS = 20_000  # cap for read_file output so large files don't blow the prompt
MAX_LIST_ENTRIES = 500

llm = ChatOpenAI(
    model="openai.gpt-oss-120b"
)

# ---------- agent workdir & path confinement ----------
# Everything the agent's read_file/write_file/list_directory/create_artifact
# tools touch lives under a "workdir". This is a *separate* tree from the
# app's own storage (DB/, tool_envs/, tools/, logs/), so a careless or
# adversarial write from agent (or LLM-generated tool) code can't reach the
# checkpoint db, sandbox venvs, or persisted tool source.
#
# DEFAULT_AGENT_WORKDIR is the fallback used when a thread hasn't picked a
# folder of its own. Cowork-style "select a working directory" support lets
# each thread point its tools at a *different* folder on disk instead --
# see DynamicAgentManager.set_working_directory() below. The active workdir
# for whichever thread is currently running is tracked in `_workdir_ctx`
# (same contextvar pattern as `_token_usage_ctx`), so every tool call made
# during that run resolves paths against the right folder without needing
# the thread_id threaded through every function signature.
DEFAULT_AGENT_WORKDIR = Path(os.environ.get("AGENT_WORKDIR", os.path.join(os.getcwd(), "agent_workspace"))).resolve()
DEFAULT_AGENT_WORKDIR.mkdir(parents=True, exist_ok=True)

_workdir_ctx: contextvars.ContextVar[Path | None] = contextvars.ContextVar("workdir_ctx", default=None)


def current_workdir() -> Path:
    """The workdir tools should confine themselves to for the run currently
    in progress: the thread's selected cowork folder if one is set, else
    DEFAULT_AGENT_WORKDIR."""
    return _workdir_ctx.get() or DEFAULT_AGENT_WORKDIR


def current_artifacts_dir() -> Path:
    """Artifacts live inside a `.artifacts` folder of whichever workdir is
    active, so a report the agent creates while cowork'd into a user's
    folder actually shows up there, rather than always landing in the
    default workspace."""
    d = current_workdir() / ".artifacts"
    d.mkdir(parents=True, exist_ok=True)
    return d


# Internal app directories the agent must never read/write/list, even if a
# workdir is ever misconfigured to overlap the project root. Resolved
# lazily (not at import time) so they reflect wherever the process actually
# ends up creating them (relative to cwd, same as the DB/tool_envs/logs
# creation calls elsewhere in this file).
_AGENT_DATA_DIR_NAMES = ("DB", "tool_envs", "tools", "logs")


def _agent_data_roots() -> list[Path]:
    return [Path(name).resolve() for name in _AGENT_DATA_DIR_NAMES]


class PathConfinementError(Exception):
    """Raised when a tool-requested path would escape the active workdir."""


class WorkdirSelectionError(Exception):
    """Raised when a user-requested working directory is invalid or refused."""


def resolve_and_confine(path_str: str | None, base: Path = None) -> Path:
    """Canonicalize a path relative to `base` (default: the active workdir)
    and reject any escape (via '..' or a symlink) outside it."""
    base = (base or current_workdir()).resolve()
    candidate = (base / path_str) if path_str else base
    resolved = candidate.resolve()
    if not (resolved == base or resolved.is_relative_to(base)):
        raise PathConfinementError(f"Path escapes the agent workdir ({base}): {path_str!r}")
    return resolved


def _path_touches_agent_data(path: Path) -> bool:
    return any(path == root or path.is_relative_to(root) for root in _agent_data_roots())


_APP_DATA_ERROR = (
    "Refused: this path is inside the application's internal data directory "
    "(DB, tool_envs, tools, logs) -- never a target for agent file operations."
)

def _as_text(content) -> str:
    """Extract plain text from a BaseMessage.content, which is normally a
    str but becomes a list of content blocks (text + image_url) once
    multimodal attachments are involved -- every place that treats
    message content as a plain string (planner goal, evaluator output,
    conversation history) needs this instead of assuming str."""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = [b.get("text", "") for b in content if isinstance(b, dict) and b.get("type") == "text"]
        return "\n".join(p for p in parts if p)
    return str(content) if content else ""


def _extract_image_blocks(content) -> list[dict]:
    """Pull the image_url content blocks (if any) out of a message's
    content, so they can be re-attached to whichever task actually needs
    to see them, instead of being lost once the goal is reduced to text."""
    if isinstance(content, list):
        return [b for b in content if isinstance(b, dict) and b.get("type") == "image_url"]
    return []


def parse_json_safely(text: str, default=None):
    """Extract and parse JSON from an LLM response, tolerating markdown fences."""
    if text is None:
        return default
    cleaned = text.strip()
    try:
        if "```json" in cleaned:
            cleaned = cleaned.split("```json")[1].split("```")[0]
        elif "```" in cleaned:
            cleaned = cleaned.split("```")[1].split("```")[0]
        return json.loads(cleaned.strip())
    except Exception:
        return default

_token_usage_ctx: contextvars.ContextVar[list | None] = contextvars.ContextVar(
    "token_usage_ctx", default=None
)


def _extract_usage(response) -> dict:
    """Pull input/output/total token counts off an LLM response in a provider-agnostic way."""
    usage = getattr(response, "usage_metadata", None)
    if usage:
        input_tokens = usage.get("input_tokens", 0)
        output_tokens = usage.get("output_tokens", 0)
        total_tokens = usage.get("total_tokens", input_tokens + output_tokens)
        return {"input_tokens": input_tokens, "output_tokens": output_tokens, "total_tokens": total_tokens}

    # Fallback: raw OpenAI-style token_usage dict inside response_metadata
    meta = getattr(response, "response_metadata", {}) or {}
    token_usage = meta.get("token_usage") or {}
    input_tokens = token_usage.get("prompt_tokens", 0)
    output_tokens = token_usage.get("completion_tokens", 0)
    total_tokens = token_usage.get("total_tokens", input_tokens + output_tokens)
    return {"input_tokens": input_tokens, "output_tokens": output_tokens, "total_tokens": total_tokens}


def _record_token_usage(node: str, response, tags: list[str] | None = None):
    """Record one LLM call's token usage against the currently-running `run()` call, if tracking is active."""
    bucket = _token_usage_ctx.get()
    if bucket is None:
        return  # no active run() context (e.g. called outside DynamicAgentManager.run)
    usage = _extract_usage(response)
    bucket.append({"node": node, "tags": tags or [], **usage})


class OrchestratorState(TypedDict):
    messages: Annotated[list[BaseMessage], add_messages]
    goal: str
    task_plan: list[dict]
    current_task_idx: int
    task_results: dict
    task_messages: list[BaseMessage]
    retry_count: int
    last_verdict: dict
    conversation_history: Annotated[list[dict], operator.add]
    attachments: list[dict]
    """image_url content blocks carried over from the triggering user
    message (see _extract_image_blocks) so any task in the plan can see
    an attached image, not just whichever task happens to read
    state["messages"] directly."""


MAX_HISTORY_TURNS_IN_PROMPT = 12  # cap how many past turns we inject into prompts


def _format_conversation_history(history: list[dict], max_turns: int = MAX_HISTORY_TURNS_IN_PROMPT) -> str:
    """Render prior conversation turns as plain text for inclusion in an LLM prompt."""
    if not history:
        return ""
    recent = history[-max_turns:]
    lines = [f"{turn.get('role', 'user').capitalize()}: {turn.get('content', '')}" for turn in recent]
    return "\n".join(lines)


# Generic runner that gets executed with the shared venv's interpreter, in a
# clean subprocess, for every dynamically generated tool call. Kept
# dependency-free (stdlib only) so it works before any pip install happens.
_SANDBOX_RUNNER_SOURCE = '''
import sys, json, importlib.util

def main():
    module_path, func_name = sys.argv[1], sys.argv[2]
    kwargs = json.loads(sys.stdin.read() or "{}")

    spec = importlib.util.spec_from_file_location("generated_tool_module", module_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    target = getattr(module, func_name)
    result = target(**kwargs)

    print(json.dumps({"ok": True, "result": result}, default=str))

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"{type(e).__name__}: {e}"}))
        sys.exit(1)
'''


class ToolSandboxExecutor:
    """
    Runs every dynamically created tool's generated code through ONE shared
    venv + a subprocess per call, so tool code can `import` any pip package
    without the main process ever needing that package (or being able to be
    crashed/hung by broken generated code).

    Trade-off vs a venv per tool: one shared site-packages means two tools
    that need conflicting versions of the same package can step on each
    other. Acceptable for an initial build where AUTO_TOOL_LIMIT already
    caps tool creation; if that ever bites, give just the conflicting tool
    its own venv (base_dir=f"tool_envs/{tool_name}") rather than switching
    the whole system over.
    """

    def __init__(self, base_dir: str = "tool_envs/shared", timeout_s: int = 25, install_timeout_s: int = 180):
        self.base_dir = base_dir
        self.timeout_s = timeout_s
        self.install_timeout_s = install_timeout_s
        self.tools_dir = os.path.join(base_dir, "tools")
        os.makedirs(self.tools_dir, exist_ok=True)

        self._runner_path = os.path.join(base_dir, "_sandbox_runner.py")
        with open(self._runner_path, "w", encoding="utf-8") as f:
            f.write(_SANDBOX_RUNNER_SOURCE)

        self._venv_dir = os.path.join(base_dir, "venv")
        self._installed_cache_path = os.path.join(base_dir, "_installed.json")
        self._installed: set[str] = set()
        self._load_installed_cache()

    # ---------- paths ----------
    def _venv_python(self) -> str:
        rel = ("Scripts", "python.exe") if os.name == "nt" else ("bin", "python")
        return os.path.join(self._venv_dir, *rel)

    def _module_path(self, tool_name: str) -> str:
        return os.path.join(self.tools_dir, f"{tool_name}.py")

    def _load_installed_cache(self):
        if os.path.exists(self._installed_cache_path):
            try:
                with open(self._installed_cache_path, "r", encoding="utf-8") as f:
                    self._installed = set(json.load(f))
            except (json.JSONDecodeError, OSError):
                self._installed = set()

    def _save_installed_cache(self):
        with open(self._installed_cache_path, "w", encoding="utf-8") as f:
            json.dump(sorted(self._installed), f)

    # ---------- setup (called once per tool, at creation time) ----------
    def ensure_env(self, tool_name: str, requirements: list[str] | None = None):
        """Create the shared venv on first use, then install only the packages not already present."""
        if not os.path.exists(self._venv_dir):
            print("🐍 Creating shared tool venv (first dynamic tool)...")
            venv.EnvBuilder(with_pip=True, clear=True).create(self._venv_dir)

        missing = [r for r in (requirements or []) if r.lower() not in self._installed]
        if missing:
            print(f"📦 Installing {missing} (needed by '{tool_name}') into shared venv...")
            proc = subprocess.run(
                [self._venv_python(), "-m", "pip", "install", "-q",
                 "--disable-pip-version-check", *missing],
                capture_output=True, text=True, timeout=self.install_timeout_s,
            )
            if proc.returncode != 0:
                raise RuntimeError(
                    f"Failed to install requirements for tool '{tool_name}': {proc.stderr[-800:]}"
                )
            self._installed.update(r.lower() for r in missing)
            self._save_installed_cache()
        elif requirements:
            print(f"📦 All requirements for '{tool_name}' already present in shared venv, skipping install.")

    def save_tool_module(self, tool_name: str, code: str) -> str:
        """Persist the generated (import-allowed) source into the shared tools folder."""
        path = self._module_path(tool_name)
        with open(path, "w", encoding="utf-8") as f:
            f.write(code)
        return path

    # ---------- execution (called on every tool invocation) ----------
    def run(self, tool_name: str, func_name: str, kwargs: dict) -> Any:
        module_path = self._module_path(tool_name)
        try:
            proc = subprocess.run(
                [self._venv_python(), self._runner_path, module_path, func_name],
                input=json.dumps(kwargs),
                capture_output=True, text=True, timeout=self.timeout_s,
            )
        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Tool '{tool_name}' timed out after {self.timeout_s}s")

        if not proc.stdout.strip():
            raise RuntimeError(f"Tool '{tool_name}' produced no output. stderr: {proc.stderr[-800:]}")

        try:
            payload = json.loads(proc.stdout.strip().splitlines()[-1])
        except json.JSONDecodeError:
            raise RuntimeError(f"Tool '{tool_name}' returned non-JSON output: {proc.stdout[-800:]}")

        if not payload.get("ok"):
            raise RuntimeError(f"Tool '{tool_name}' error: {payload.get('error')}")
        return payload["result"]


# ---------- artifact content converters (used by create_artifact tool) ----------
_ARTIFACT_KINDS = {"doc", "pdf", "slides", "image"}

_INLINE_MD_RE = re.compile(
    r"\*\*(?P<bold1>.+?)\*\*"
    r"|__(?P<bold2>.+?)__"
    r"|\*(?P<italic1>.+?)\*"
    r"|_(?P<italic2>.+?)_"
    r"|`(?P<code>.+?)`"
)


def _add_inline_runs(paragraph, text: str) -> None:
    """Splits a line on inline markdown (bold/italic/code) into styled runs
    so a downloaded .docx shows real bold/italic instead of literal '**x**'."""
    pos = 0
    for match in _INLINE_MD_RE.finditer(text):
        if match.start() > pos:
            paragraph.add_run(text[pos:match.start()])
        if match.group("bold1") is not None:
            paragraph.add_run(match.group("bold1")).bold = True
        elif match.group("bold2") is not None:
            paragraph.add_run(match.group("bold2")).bold = True
        elif match.group("italic1") is not None:
            paragraph.add_run(match.group("italic1")).italic = True
        elif match.group("italic2") is not None:
            paragraph.add_run(match.group("italic2")).italic = True
        elif match.group("code") is not None:
            paragraph.add_run(match.group("code")).font.name = "Courier New"
        pos = match.end()
    if pos < len(text) or pos == 0:
        paragraph.add_run(text[pos:])


def _parse_md_table_row(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _is_md_table_separator(line: str) -> bool:
    cells = _parse_md_table_row(line)
    return bool(cells) and all(re.fullmatch(r":?-+:?", c) for c in cells)


def _add_docx_table(document, header: list[str], rows: list[list[str]]) -> None:
    table = document.add_table(rows=1 + len(rows), cols=len(header))
    table.style = "Table Grid"
    for col, text in enumerate(header):
        paragraph = table.rows[0].cells[col].paragraphs[0]
        _add_inline_runs(paragraph, text)
        for run in paragraph.runs:
            run.bold = True
    for r, row in enumerate(rows, start=1):
        for col, text in enumerate(row):
            if col < len(header):
                _add_inline_runs(table.rows[r].cells[col].paragraphs[0], text)


def _md_to_docx(content: str):
    """Small, deliberately non-exhaustive Markdown -> docx converter: enough
    structure for agent-authored reports (headings, bullets, paragraphs,
    tables, inline bold/italic/code), not a full CommonMark implementation."""
    document = _DocxDocument()
    lines = content.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("|") and i + 1 < len(lines) and _is_md_table_separator(lines[i + 1]):
            header = _parse_md_table_row(stripped)
            i += 2
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(_parse_md_table_row(lines[i]))
                i += 1
            _add_docx_table(document, header, rows)
            continue
        if stripped.startswith("### "):
            _add_inline_runs(document.add_heading("", level=3), stripped[4:])
        elif stripped.startswith("## "):
            _add_inline_runs(document.add_heading("", level=2), stripped[3:])
        elif stripped.startswith("# "):
            _add_inline_runs(document.add_heading("", level=1), stripped[2:])
        elif stripped.startswith(("- ", "* ")):
            _add_inline_runs(document.add_paragraph(style="List Bullet"), stripped[2:])
        else:
            _add_inline_runs(document.add_paragraph(), stripped)
        i += 1
    return document


def _build_pptx(slides: list[dict]):
    presentation = _PptxPresentation()
    layout = presentation.slide_layouts[1]  # "Title and Content"
    for entry in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = entry["title"]
        bullets = entry.get("bullets") or []
        if bullets:
            body = slide.placeholders[1].text_frame
            body.text = bullets[0]
            for bullet in bullets[1:]:
                body.add_paragraph().text = bullet
    return presentation


class DynamicToolRegistry:
    """Manages dynamic tool creation and registration."""
    def __init__(self):
        self.tools: dict[str, BaseTool] = {}
        self.tool_code: dict[str, str] = {}
        self.artifacts: list[dict] = []
        """Side-channel that create_artifact appends to on success, mirroring
        cowork's ctx.artifacts pattern -- lets the caller (Streamlit, etc.)
        list what was produced during a run without parsing tool output."""
        self.sandbox = ToolSandboxExecutor()
        self._register_default_tools()

    def _register_default_tools(self):
        """Register built-in tools"""
        self.register_tool("search", self._search_tool())
        self.register_tool("calculator", self._calculator_tool())
        self.register_tool("stock_price", self._stock_price_tool())
        self.register_tool("read_document", self._read_document_tool())
        self.register_tool("list_documents", self._list_documents_tool())
        self.register_tool("generate_document", self._generate_document_tool())
        self.register_tool("read_file", self._read_file_tool())
        self.register_tool("write_file", self._write_file_tool())
        self.register_tool("list_directory", self._list_directory_tool())
        self.register_tool("view_image", self._view_image_tool())
        self.register_tool("create_artifact", self._create_artifact_tool())

    def list_artifacts(self) -> list[dict]:
        return list(self.artifacts)

    def register_tool(self, name: str, tool_obj: BaseTool, code: str = ""):
        """Register a tool, forcing tool_obj.name to match the registry key."""
        tool_obj.name = name
        self.tools[name] = tool_obj
        self.tool_code[name] = code
        if code:
            try:
                tools_dir = os.path.join(os.path.dirname(__file__), "tools")
                os.makedirs(tools_dir, exist_ok=True)
                with open(os.path.join(tools_dir, f"{name}.py"), "w", encoding="utf-8") as f:
                    f.write(f"# Auto-generated tool: {name}\n# Created: {datetime.now().isoformat()}\n\n{code}")
            except OSError as e:
                print(f"⚠️ Could not persist source for tool '{name}' to disk: {e}")
        print(f"✅ Tool registered: {name}")

    def get_all_tools(self) -> list[BaseTool]:
        """Get all registered tools as a list"""
        return list(self.tools.values())

    def get_tool(self, name: str) -> BaseTool | None:
        """Get a specific tool by name"""
        return self.tools.get(name)

    def list_tools(self) -> dict[str, str]:
        """List all tools with their descriptions"""
        return {name: tool.description for name, tool in self.tools.items()}

    _DEFAULT_TOOL_NAMES = (
        "search", "calculator", "stock_price",
        "read_document", "list_documents", "generate_document",
        "read_file", "write_file", "list_directory", "view_image", "create_artifact",
    )

    def remove_tool(self, name: str, delete_generated_code: bool = True) -> bool:
        """
        Remove a dynamically created tool so it's no longer offered to any
        agent. Deliberately does NOT touch the shared venv's installed
        packages or the install cache — those stay put so that if this tool
        (or another one needing the same packages) is auto-created again
        later, it won't need to reinstall anything.

        delete_generated_code=False keeps the .py source files on disk
        (e.g. if you just want to unregister it temporarily but might
        restore the exact same code later); default True removes them so a
        future creation for the same tool_name starts clean.
        """
        if name not in self.tools:
            print(f"⚠️ Tool '{name}' not found in registry, nothing to remove.")
            return False
        if name in self._DEFAULT_TOOL_NAMES:
            print(f"⚠️ Refusing to remove built-in tool '{name}'.")
            return False

        del self.tools[name]
        self.tool_code.pop(name, None)

        if delete_generated_code:
            # the copy register_tool() persisted for reference/inspection
            persisted_path = os.path.join(os.path.dirname(__file__), "tools", f"{name}.py")
            if os.path.exists(persisted_path):
                os.remove(persisted_path)
            # the sandbox's own copy, so create_tool_from_prompt starts fresh next time
            module_path = self.sandbox._module_path(name)
            if os.path.exists(module_path):
                os.remove(module_path)

        print(f"🗑️ Tool '{name}' removed from registry.")
        return True

    @staticmethod
    def _search_tool() -> BaseTool:
        @tool
        def search(query: str) -> str:
            """Search the web for information"""
            search_tool = GoogleSerperAPIWrapper(region="us-en")
            result = search_tool.run(query)
            return result
        return search

    @staticmethod
    def _calculator_tool() -> BaseTool:
        @tool
        def calculator(first_num: float, second_num: float, operation: str) -> dict:
            """Perform arithmetic: add, sub, mul, div"""
            ops = {
                "add": lambda a, b: a + b,
                "sub": lambda a, b: a - b,
                "mul": lambda a, b: a * b,
                "div": lambda a, b: a / b if b != 0 else "Error: Division by zero"
            }
            result = ops.get(operation, lambda a, b: "Error: Unknown operation")(first_num, second_num)
            return {"result": result, "operation": operation}
        return calculator

    @staticmethod
    def _stock_price_tool() -> BaseTool:
        @tool
        def get_stock_price(symbol: str) -> dict:
            """Fetch stock price for a symbol (e.g., AAPL, TSLA)"""
            api_key = os.environ.get("ALPHAVANTAGE_API_KEY")
            if not api_key:
                return {"error": "ALPHAVANTAGE_API_KEY not set in environment"}
            url = f"https://www.alphavantage.co/query?function=GLOBAL_QUOTE&symbol={symbol}&apikey={api_key}"
            try:
                r = requests.get(url, timeout=5)
                return r.json()
            except Exception as e:
                return {"error": str(e)}
        return get_stock_price

    @staticmethod
    def _read_document_tool() -> BaseTool:
        @tool
        def read_document(doc_id: str) -> dict:
            """Read the content of an uploaded document by its doc_id.

            Use this tool when the user asks about, references, or wants to
            read a previously uploaded file. Returns the extracted text along
            with metadata (filename, character count, upload time).

            Parameters
            ----------
            doc_id : str
                The document ID returned by the upload endpoint.
            """
            doc = docpipe.document_store.get(doc_id)
            if doc is None:
                return {"error": f"No document found with id: {doc_id}"}
            return {
                "doc_id": doc["doc_id"],
                "filename": doc["filename"],
                "char_count": doc["char_count"],
                "uploaded_at": doc["uploaded_at"],
                "content": doc["text"][:docpipe.MAX_CHARS_FOR_CONTEXT],
                "truncated": len(doc["text"]) > docpipe.MAX_CHARS_FOR_CONTEXT,
            }
        return read_document

    @staticmethod
    def _list_documents_tool() -> BaseTool:
        @tool
        def list_documents(query: str = "") -> dict:
            """List all uploaded documents with their IDs, filenames, sizes, and upload times.

            Use this tool when you need to find which documents have been
            uploaded, or when the user references a document by name and you
            need to look up its doc_id.

            Parameters
            ----------
            query : str
                Optional search filter (not yet implemented, pass empty string).
            """
            docs = docpipe.document_store.list()
            return {
                "count": len(docs),
                "documents": docs,
            }
        return list_documents

    @staticmethod
    def _generate_document_tool() -> BaseTool:
        @tool
        def generate_document(content: str, filename: str = "document", format: str = "pdf") -> dict:
            """Generate a document file (PDF, DOCX, or TXT) from text content and save it to disk.

            Use this tool when the user asks you to create, generate, or export
            a document, report, summary, or any written content as a
            downloadable file.

            The generated file is saved to a 'generated_documents' folder and
            a download URL is returned.

            Parameters
            ----------
            content : str
                The text or markdown content to write into the document.
                Supports markdown headings (#) and bullet points (- or *).
            filename : str
                Base name for the output file (e.g. 'quarterly_report').
                The correct extension is added automatically.
            format : str
                Output format: 'pdf', 'docx', or 'txt'. Defaults to 'pdf'.
            """
            return docpipe.generate_document(content, filename, format)
        return generate_document

    @staticmethod
    def _read_file_tool() -> BaseTool:
        @tool
        def read_file(path: str, start_line: int = 0, end_line: int = 0) -> dict:
            """Read a file from the current agent workspace, optionally by line
            range. Handles any file type the workspace might contain:

            - Plain text / code files: read directly, optionally sliced by
              start_line/end_line.
            - .pdf / .docx / .csv / .xlsx: extracted to plain text
              automatically (same extractor the upload pipeline uses).
            - Images (.png/.jpg/.jpeg/.gif/.webp): NOT handled here -- use
              the `view_image` tool instead, which attaches the image
              itself for you to see rather than returning text.

            Parameters
            ----------
            path : str
                Path relative to the current working directory. Must not be
                absolute or contain '..' -- both are rejected.
            start_line : int
                1-indexed first line to include (plain-text files only). 0
                (default) means "from the start".
            end_line : int
                1-indexed last line to include (plain-text files only). 0
                (default) means "to the end".
            """
            try:
                resolved = resolve_and_confine(path)
            except PathConfinementError as e:
                return {"error": str(e)}
            if _path_touches_agent_data(resolved):
                return {"error": _APP_DATA_ERROR}
            if not resolved.is_file():
                return {"error": f"Not a file: {path}"}

            ext = resolved.suffix.lower()
            if docpipe.is_image_extension(ext):
                return {
                    "error": (
                        f"'{path}' is an image file. Use the `view_image` tool to "
                        "look at it instead of read_file."
                    )
                }

            if ext in docpipe.TEXT_EXTENSIONS - {".txt"}:
                # pdf / docx / csv / xlsx -- extract to text via the same
                # pipeline uploads go through, instead of dumping raw bytes.
                try:
                    content = docpipe.extract_text(resolved.name, resolved.read_bytes())
                except Exception as e:
                    return {"error": f"Could not extract text from '{path}': {type(e).__name__}: {e}"}
                truncated = len(content) > MAX_READ_CHARS
                if truncated:
                    content = content[:MAX_READ_CHARS]
                return {"path": path, "content": content, "truncated": truncated, "extracted_from": ext}

            lines = resolved.read_text(errors="replace").splitlines(keepends=True)
            if start_line or end_line:
                start_idx = max((start_line or 1) - 1, 0)
                end_idx = end_line if end_line else len(lines)
                lines = lines[start_idx:end_idx]
            content = "".join(lines)
            truncated = len(content) > MAX_READ_CHARS
            if truncated:
                content = content[:MAX_READ_CHARS]
            return {"path": path, "content": content, "truncated": truncated}
        return read_file

    @staticmethod
    def _view_image_tool() -> BaseTool:
        @tool
        def view_image(path: str) -> dict:
            """Load an image file from the current agent workspace so you can
            actually see it. Use this for .png/.jpg/.jpeg/.gif/.webp files
            instead of read_file -- read_file only handles text-representable
            content and will refuse image paths.

            After this tool runs, the image is attached to the conversation
            for you to look at directly on your next turn.

            Parameters
            ----------
            path : str
                Path relative to the current working directory. Must not be
                absolute or contain '..' -- both are rejected.
            """
            try:
                resolved = resolve_and_confine(path)
            except PathConfinementError as e:
                return {"error": str(e)}
            if _path_touches_agent_data(resolved):
                return {"error": _APP_DATA_ERROR}
            if not resolved.is_file():
                return {"error": f"Not a file: {path}"}

            ext = resolved.suffix.lower()
            if not docpipe.is_image_extension(ext):
                return {"error": f"'{path}' is not a supported image type ({sorted(docpipe.IMAGE_EXTENSIONS)})."}

            data = resolved.read_bytes()
            if len(data) > docpipe.MAX_IMAGE_BYTES:
                return {
                    "error": (
                        f"Image is {len(data)} bytes, over the {docpipe.MAX_IMAGE_BYTES}-byte limit "
                        "for viewing."
                    )
                }

            import base64
            return {
                "is_image": True,
                "path": path,
                "mime_type": docpipe.image_mime_type(ext),
                "data_b64": base64.b64encode(data).decode("ascii"),
            }
        return view_image

    @staticmethod
    def _write_file_tool() -> BaseTool:
        @tool
        def write_file(path: str, content: str) -> dict:
            """Create or overwrite a file in the agent workspace with the given
            text content, creating parent directories as needed.

            Parameters
            ----------
            path : str
                Path relative to the agent workspace. Must not be absolute
                or contain '..' -- both are rejected.
            content : str
                Full text content to write. This REPLACES the file if it
                already exists.
            """
            try:
                resolved = resolve_and_confine(path)
            except PathConfinementError as e:
                return {"error": str(e)}
            if _path_touches_agent_data(resolved):
                return {"error": _APP_DATA_ERROR}
            resolved.parent.mkdir(parents=True, exist_ok=True)
            resolved.write_text(content, encoding="utf-8")
            return {"status": "ok", "path": path, "bytes_written": len(content.encode("utf-8"))}
        return write_file

    @staticmethod
    def _list_directory_tool() -> BaseTool:
        @tool
        def list_directory(path: str = "") -> dict:
            """List files and directories at a path inside the agent workspace
            (non-recursive). Pass an empty string for the workspace root.

            Parameters
            ----------
            path : str
                Path relative to the agent workspace. Must not be absolute
                or contain '..' -- both are rejected.
            """
            try:
                resolved = resolve_and_confine(path or None)
            except PathConfinementError as e:
                return {"error": str(e)}
            if _path_touches_agent_data(resolved):
                return {"error": _APP_DATA_ERROR}
            if not resolved.is_dir():
                return {"error": f"Not a directory: {path}"}

            entries = sorted(resolved.iterdir(), key=lambda p: p.name)[:MAX_LIST_ENTRIES]
            return {
                "path": path,
                "entries": [f"{'d' if e.is_dir() else 'f'} {e.name}" for e in entries],
            }
        return list_directory

    def _create_artifact_tool(self) -> BaseTool:
        """Instance method (not static) so the handler can append to
        self.artifacts -- everything else in this registry is a static
        closure, but create_artifact needs a side-channel back to the
        registry, same reasoning as cowork's ctx.artifacts."""
        registry = self

        @tool
        def create_artifact(filename: str, kind: str, title: str, content: str) -> dict:
            """Create a distinct, previewable deliverable (report, diagram, or
            slide deck) in the agent workspace's .artifacts folder, as opposed
            to an ordinary working file written via write_file.

            Parameters
            ----------
            filename : str
                Base filename without extension.
            kind : str
                One of: "doc" (Markdown -> .md + .docx), "pdf" (Markdown ->
                .pdf), "slides" (JSON -> .pptx), "image" (raw SVG -> .svg).
            title : str
                Human-readable title for the artifact.
            content : str
                Shape depends on kind:
                - doc / pdf: Markdown text (headings, bullets, tables,
                  inline bold/italic/code all supported).
                - slides: JSON string like
                  {"slides": [{"title": "...", "bullets": ["...", "..."]}]}
                - image: raw SVG markup starting with "<svg" or "<?xml".
            """
            if kind not in _ARTIFACT_KINDS:
                return {"error": f"Unknown artifact kind: {kind!r}. Must be one of {sorted(_ARTIFACT_KINDS)}."}

            artifact_id = uuid.uuid4().hex
            stem = f"{artifact_id}__{filename}"
            out_dir = current_artifacts_dir()

            try:
                if kind == "doc":
                    if _DocxDocument is None:
                        return {"error": "python-docx is not installed. Run: pip install python-docx"}
                    (out_dir / f"{stem}.md").write_text(content, encoding="utf-8")
                    _md_to_docx(content).save(out_dir / f"{stem}.docx")
                    preview_path = out_dir / f"{stem}.md"
                    download_path = out_dir / f"{stem}.docx"

                elif kind == "pdf":
                    if _weasyprint is None or _markdown is None:
                        return {"error": "weasyprint/markdown are not installed. Run: pip install weasyprint markdown"}
                    html = _markdown.markdown(content, extensions=["extra"])
                    pdf_path = out_dir / f"{stem}.pdf"
                    _weasyprint.HTML(string=html).write_pdf(str(pdf_path))
                    preview_path = download_path = pdf_path

                elif kind == "slides":
                    if _PptxPresentation is None:
                        return {"error": "python-pptx is not installed. Run: pip install python-pptx"}
                    try:
                        parsed = json.loads(content)
                        slides = parsed["slides"]
                        if not isinstance(slides, list) or not slides:
                            raise ValueError("'slides' must be a non-empty array")
                        for entry in slides:
                            if not isinstance(entry.get("title"), str):
                                raise ValueError("each slide needs a string 'title'")
                    except (json.JSONDecodeError, KeyError, TypeError, ValueError) as exc:
                        return {
                            "error": (
                                "Invalid slides content -- expected JSON like "
                                '{"slides": [{"title": "...", "bullets": ["...", "..."]}]}. '
                                f"Error: {exc}"
                            )
                        }
                    (out_dir / f"{stem}.json").write_text(json.dumps(parsed), encoding="utf-8")
                    _build_pptx(slides).save(out_dir / f"{stem}.pptx")
                    preview_path = out_dir / f"{stem}.json"
                    download_path = out_dir / f"{stem}.pptx"

                else:  # image
                    stripped = content.strip()
                    if not (stripped.startswith("<svg") or stripped.startswith("<?xml")):
                        return {"error": "Image artifact content must be raw SVG markup starting with <svg ...> or <?xml ...>."}
                    svg_path = out_dir / f"{stem}.svg"
                    svg_path.write_text(content, encoding="utf-8")
                    preview_path = download_path = svg_path
            except Exception as e:
                return {"error": f"Failed to build {kind} artifact: {type(e).__name__}: {e}"}

            record = {
                "artifact_id": artifact_id,
                "filename": filename,
                "kind": kind,
                "title": title,
                "preview_path": str(preview_path),
                "download_path": str(download_path),
            }
            registry.artifacts.append(record)
            return {"status": "created", **record}
        return create_artifact

    @traceable(name="create_tool_from_prompt", run_type="chain")
    def create_tool_from_prompt(self, prompt: str, tool_name: str) -> BaseTool:
        """Create a new tool from a natural language prompt, run in an isolated venv + subprocess sandbox."""
        print(f"🔨 Creating tool '{tool_name}' from prompt...")

        creation_prompt = f"""
        Create a Python tool for the following requirement:
        {prompt}

        This code will run in its own isolated subprocess (its own venv), so
        you MAY use import statements for third-party packages — anything on
        PyPI is fine. Do NOT use os, subprocess, eval, or exec.

        Respond with ONLY a single JSON object, no prose, no markdown fences,
        in exactly this shape:
        {{
          "requirements": ["package_name", ...],
          "function_name": "{tool_name}",
          "docstring": "one sentence describing what this tool does, shown to the agent as the tool's description",
          "args_schema": {{
            "arg_name": {{"type": "string", "description": "..."}}
          }},
          "code": "def {tool_name}(arg_name: str) -> dict:\\n    ...\\n    return {{...}}"
        }}

        Rules for "requirements": pip package names this code imports; [] if
        none beyond the standard library.
        Rules for "code":
        1. A single top-level function named exactly "{tool_name}" whose
           parameters match args_schema.
        2. No @tool decorator — this file is imported and called directly
           inside the sandbox, not through langchain.
        3. Handle errors gracefully (try/except) and always return a
           JSON-serializable value.
        4. Keep it focused and simple.
        5. If the task requires returning file/image bytes, base64-encode
           them (e.g. base64.b64encode(data).decode("utf-8")) as a string
           field rather than returning only an external URL.
        6. Valid "type" values for args_schema entries: string, integer,
           number, boolean, array, object.
        """

        response = llm.invoke(
            creation_prompt,
            config={
                "run_name": "tool_codegen_llm",
                "tags": ["tool_creation", tool_name],
                "metadata": {"tool_name": tool_name},
            },
        )
        _record_token_usage(f"tool_codegen:{tool_name}", response, ["tool_creation"])
        current_spec_text = response.content
        print(f"📝 Generated spec for tool '{tool_name}':\n{current_spec_text}")

        last_error = None

        for attempt in range(MAX_TOOL_VALIDATION_RETRIES + 1):
            try:
                spec = parse_json_safely(self._extract_code_block(current_spec_text), default=None)
                if not spec or "code" not in spec or "function_name" not in spec:
                    raise ValueError("Model did not return the expected JSON tool spec (missing 'code' or 'function_name').")

                cleaned_code = spec["code"]
                requirements = spec.get("requirements") or []

                self.sandbox.ensure_env(tool_name, requirements)
                self.sandbox.save_tool_module(tool_name, cleaned_code)

                created_tool = self._build_structured_tool(
                    tool_name=tool_name,
                    function_name=spec["function_name"],
                    docstring=spec.get("docstring", f"Dynamically created tool: {tool_name}"),
                    args_schema=spec.get("args_schema", {}),
                )

                ok, smoke_message = self._smoke_test_tool(created_tool)
                if not ok:
                    raise RuntimeError(f"Smoke test failed: {smoke_message}")

                self.register_tool(tool_name, created_tool, cleaned_code)
                print(f"✅ Tool '{tool_name}' created and passed smoke test.")
                return created_tool

            except Exception as e:
                last_error = str(e)
                print(f"⚠️ Tool '{tool_name}' generation attempt {attempt + 1} failed: {last_error}")

                if attempt >= MAX_TOOL_VALIDATION_RETRIES:
                    break

                repair_prompt = f"""
                    The following tool spec failed with this error:
                    {last_error}

                    Original spec:
                    {current_spec_text}

                    Fix it so it runs correctly and satisfies the original requirement:
                    {prompt}

                    Same rules as before: respond with ONLY the corrected JSON
                    object (requirements, function_name, docstring,
                    args_schema, code). The function may use import
                    statements for third-party packages, but not os,
                    subprocess, eval, or exec.
                """
                response = llm.invoke(
                    repair_prompt,
                    config={
                        "run_name": "tool_codegen_repair_llm",
                        "tags": ["tool_creation", "repair", tool_name],
                        "metadata": {"tool_name": tool_name, "error": last_error},
                    },
                )
                _record_token_usage(f"tool_codegen_repair:{tool_name}", response, ["tool_creation", "repair"])
                current_spec_text = response.content

        print(f"❌ Error creating tool '{tool_name}': {last_error}")
        raise RuntimeError(
            f"Failed to create a working tool '{tool_name}' after "
            f"{MAX_TOOL_VALIDATION_RETRIES + 1} attempt(s). Last error: {last_error}"
        )

    @staticmethod
    def _extract_code_block(raw: str) -> str:
        if "```json" in raw:
            return raw.split("```json")[1].split("```")[0]
        elif "```python" in raw:
            return raw.split("```python")[1].split("```")[0]
        elif "```" in raw:
            return raw.split("```")[1].split("```")[0]
        return raw

    @staticmethod
    def _args_model(tool_name: str, args_schema: dict):
        """Turn the LLM's args_schema JSON into a pydantic model for StructuredTool."""
        if not args_schema:
            return None
        type_map = {
            "string": str, "integer": int, "number": float,
            "boolean": bool, "array": list, "object": dict,
        }
        fields = {}
        for arg_name, info in args_schema.items():
            py_type = type_map.get((info or {}).get("type", "string"), str)
            description = (info or {}).get("description", "")
            fields[arg_name] = (py_type, Field(..., description=description))
        return create_model(f"{tool_name}_Args", **fields)

    def _build_structured_tool(self, tool_name: str, function_name: str, docstring: str, args_schema: dict) -> BaseTool:
        """Wrap a call into the sandbox as a normal langchain BaseTool, so nothing downstream changes."""
        args_model = self._args_model(tool_name, args_schema)

        def _invoke(**kwargs):
            return self.sandbox.run(tool_name, function_name, kwargs)

        return StructuredTool.from_function(
            func=_invoke,
            name=tool_name,
            description=docstring,
            args_schema=args_model,
        )

    @staticmethod
    def _generate_dummy_args(tool_obj: BaseTool) -> dict:
        """Build plausible synthetic args from a tool's schema for smoke testing."""
        dummy = {}
        try:
            schema_props = tool_obj.args or {}
        except Exception:
            schema_props = {}
        for field_name, field_info in schema_props.items():
            field_type = (field_info or {}).get("type", "string")
            if field_type == "integer":
                dummy[field_name] = 1
            elif field_type == "number":
                dummy[field_name] = 1.0
            elif field_type == "boolean":
                dummy[field_name] = True
            elif field_type == "array":
                dummy[field_name] = []
            elif field_type == "object":
                dummy[field_name] = {}
            else:
                dummy[field_name] = "test"
        return dummy

    def _smoke_test_tool(self, tool_obj: BaseTool) -> tuple[bool, str]:
        """Invoke a freshly generated tool with synthetic args to catch obvious runtime errors."""
        dummy_args = self._generate_dummy_args(tool_obj)
        try:
            result = tool_obj.invoke(dummy_args)
            json.dumps(result, default=str)
            return True, "ok"
        except Exception as e:
            return False, f"{type(e).__name__}: {e}"


class DynamicAgentFactory:
    """Creates and caches specialized sub-agents (system_prompt, tools, llm) at runtime."""
    def __init__(self, tool_registry: DynamicToolRegistry):
        self.tool_registry = tool_registry
        self.agents: dict[str, dict] = {}

    def get_or_create(self, role: str, task_description: str) -> dict:
        if role in self.agents:
            return self.refresh_tools(role, task_description)
        return self.create_agent(role, task_description)

    @traceable(name="auto_create_missing_tools", run_type="chain")
    def _create_missing_tools(self, specs: list[dict]) -> list[str]:
        """Auto-create tools an LLM tool-selection step flagged as missing, capped at AUTO_TOOL_LIMIT."""
        created = []
        for spec in (specs or [])[:AUTO_TOOL_LIMIT]:
            name = (spec or {}).get("name")
            prompt = (spec or {}).get("prompt")
            if not name or not prompt:
                continue
            if self.tool_registry.get_tool(name) is not None:
                created.append(name)
                continue
            try:
                print(f"⚡ Auto-creating missing tool '{name}' from agent workflow...")
                self.tool_registry.create_tool_from_prompt(prompt, name)
                created.append(name)
            except Exception as e:
                print(f"❌ Auto tool creation failed for '{name}': {e}")
        return created

    @traceable(name="refresh_agent_tools", run_type="chain")
    def refresh_tools(self, role: str, task_description: str) -> dict:
        """Re-check an existing agent's tool set against the current registry and task. Only expands tools, never removes."""
        agent_conf = self.agents[role]
        available_tools = self.tool_registry.list_tools()

        selection_prompt = f"""
The "{role}" agent currently has these tools: {agent_conf['tool_names']}
It now needs to handle this task: "{task_description}"

Available tools (name: description): {json.dumps(available_tools)}

Return ONLY JSON in this exact shape:
{{
  "tools": ["tool_name1", "tool_name2"],
  "new_tools": [{{"name": "snake_case_tool_name", "prompt": "one self-contained instruction describing what this new tool must do"}}]
}}

Rules:
- "tools": its current tools plus any additional EXISTING tools (from the
  Available tools list) this new task genuinely needs. Do not drop a
  currently held tool unless it is clearly irrelevant.
- "new_tools": ONLY include an entry here if none of the available tools can
  do something this task clearly requires (e.g. reading email, sending a
  message, converting a file format, calling a specific API). Leave it as
  an empty list if the available tools are sufficient.
- Do not propose more than {AUTO_TOOL_LIMIT} new tools.
"""
        response = llm.invoke(
            selection_prompt,
            config={
                "run_name": "agent_tool_refresh_llm",
                "tags": ["tool_refresh", role],
                "metadata": {"role": role, "task_description": task_description},
            },
        )
        _record_token_usage(f"tool_refresh:{role}", response, ["tool_refresh"])
        result = parse_json_safely(response.content, default=None)
        if result is None:
            print(f"⚠️ Tool refresh JSON parse failed for role '{role}'; keeping existing tools.")
            return agent_conf

        created_tool_names = self._create_missing_tools(result.get("new_tools", []))

        desired_names = (
            set(result.get("tools", []))
            | set(agent_conf["tool_names"])
            | set(created_tool_names)
        )
        if desired_names == set(agent_conf["tool_names"]):
            return agent_conf

        selected_tools = [
            self.tool_registry.get_tool(name)
            for name in desired_names
            if self.tool_registry.get_tool(name) is not None
        ]
        agent_conf["llm"] = llm.bind_tools(selected_tools) if selected_tools else llm
        agent_conf["tool_names"] = [t.name for t in selected_tools]
        print(f"🔄 Agent '{role}' tools refreshed: {agent_conf['tool_names']}")
        return agent_conf

    def tool_directive(self, tool_names: list[str]) -> str:
        """Fixed instruction appended to every agent's system prompt to force tool use when relevant."""
        if not tool_names:
            return ""
        descriptions = self.tool_registry.list_tools()
        lines = [f'- {name}: {descriptions.get(name, "")}' for name in tool_names]
        return (
            "\n\nTOOL USE GUIDELINES:\n"
            + "\n".join(lines)
            + "\n\n1. If completing this task requires external or live data (prices, search results, computation), call the appropriate tool above.\n"
            "2. CRITICAL: If you have already executed a tool call and received results in a ToolMessage, DO NOT call the tool again for the same query. Synthesize your final answer using the retrieved results immediately.\n"
            "3. Do NOT make redundant or repeated tool calls once information has been fetched."
        )

    @traceable(name="create_agent", run_type="chain")
    def create_agent(self, role: str, task_description: str) -> dict:
        _log("WORKFLOW", f"🧬 Creating new agent for role: '{role}'")
        available_tools = self.tool_registry.list_tools()

        selection_prompt = f"""
You are configuring a specialized AI agent for the role "{role}".
It will handle tasks like: "{task_description}"

Available tools (name: description): {json.dumps(available_tools)}

Return ONLY JSON in this exact shape:
{{
  "system_prompt": "a system prompt defining this agent's persona, scope and behavior",
  "tools": ["tool_name1", "tool_name2"],
  "new_tools": [{{"name": "snake_case_tool_name", "prompt": "one self-contained instruction describing what this new tool must do"}}]
}}

Rules:
- "tools": only names that already appear in the Available tools list above.
  Select tools genuinely relevant to this role. It is fine to select zero.
- "new_tools": ONLY include an entry here if none of the available tools can
  do something this role clearly needs (e.g. reading email, sending a
  message, converting a file format, calling a specific external API).
  Leave it as an empty list if the available tools are sufficient.
- Do not propose more than {AUTO_TOOL_LIMIT} new tools.
"""
        _log("AI-REQUEST", f"Invoking agent-creation LLM for role '{role}'")
        response = llm.invoke(
            selection_prompt,
            config={
                "run_name": "agent_creation_llm",
                "tags": ["agent_creation", role],
                "metadata": {"role": role, "task_description": task_description},
            },
        )
        _record_token_usage(f"agent_creation:{role}", response, ["agent_creation"])
        _log_block("AI-REPLY", f"Agent-creation reply for role '{role}'", _as_text(response.content))
        config = parse_json_safely(response.content, default=None)
        if config is None:
            _log("AI-REPLY", f"⚠️ Tool selection JSON parse failed for role '{role}'. Raw response: {response.content[:300]!r}")
            config = {
                "system_prompt": f"You are a focused, helpful '{role}' agent. Be concise and accurate.",
                "tools": [],
                "new_tools": [],
            }

        created_tool_names = self._create_missing_tools(config.get("new_tools", []))
        all_tool_names = list(dict.fromkeys(config.get("tools", []) + created_tool_names))

        selected_tools = [
            self.tool_registry.get_tool(name)
            for name in all_tool_names
            if self.tool_registry.get_tool(name) is not None
        ]

        agent_llm = llm.bind_tools(selected_tools) if selected_tools else llm

        agent_conf = {
            "role": role,
            "system_prompt": config.get("system_prompt", f"You are a helpful '{role}' agent."),
            "tool_names": [t.name for t in selected_tools],
            "llm": agent_llm,
        }
        self.agents[role] = agent_conf
        _log("WORKFLOW", f"✅ Agent '{role}' created with tools: {agent_conf['tool_names']}")
        return agent_conf

    def list_agents(self) -> dict:
        return {
            role: {"system_prompt": conf["system_prompt"], "tools": conf["tool_names"]}
            for role, conf in self.agents.items()
        }

    def remove_tool_from_agents(self, tool_name: str):
        """Strip a removed tool out of every agent that currently holds it, and
        rebind each affected agent's llm so the tool no longer appears in the
        schema it sees or in its tool_directive. Leaves agents unaffected if
        they never had this tool."""
        for role, conf in self.agents.items():
            if tool_name not in conf["tool_names"]:
                continue
            conf["tool_names"] = [t for t in conf["tool_names"] if t != tool_name]
            remaining_tools = [
                self.tool_registry.get_tool(n) for n in conf["tool_names"]
            ]
            remaining_tools = [t for t in remaining_tools if t is not None]
            conf["llm"] = llm.bind_tools(remaining_tools) if remaining_tools else llm
            print(f"🔄 Agent '{role}' tools after removing '{tool_name}': {conf['tool_names']}")


class DynamicAgentManager:
    """
    Orchestration graph:
    planner -> agent_executor <-> tools -> evaluator -> (retry | next task) -> assembler
    """
    def __init__(self):
        self.tool_registry = DynamicToolRegistry()
        self.agent_factory = DynamicAgentFactory(self.tool_registry)
        self.behavior_style = "standard"
        self.extra_instruction = ""
        self.temperature = 0.0
        self.chatbot = self._build_graph()
        self._thread_workdirs: dict[str, Path] = {}
        """Cowork-style per-thread working directory selection: each
        thread_id can point its file tools (read_file/write_file/
        list_directory/view_image/create_artifact) at a different folder
        on disk. Threads that never call set_working_directory keep using
        DEFAULT_AGENT_WORKDIR."""

    # ---------- working directory (cowork) ----------
    def set_working_directory(self, thread_id: str, path: str) -> dict:
        """Point a thread's file tools at an existing folder on disk.

        Refuses paths that don't exist, aren't directories, or land inside
        the application's own internal data directories (DB, tool_envs,
        tools, logs) -- same protection resolve_and_confine gives ordinary
        tool calls, applied up front so a bad selection fails clearly
        instead of surfacing as confusing tool errors later.
        """
        resolved = Path(path).expanduser().resolve()
        if not resolved.exists():
            raise WorkdirSelectionError(f"No such directory: {path}")
        if not resolved.is_dir():
            raise WorkdirSelectionError(f"Not a directory: {path}")
        if _path_touches_agent_data(resolved):
            raise WorkdirSelectionError(
                f"Refused: '{path}' is inside the application's internal data directory."
            )
        self._thread_workdirs[thread_id] = resolved
        print(f"📂 Thread {thread_id} working directory set to: {resolved}")
        return {"status": "ok", "thread_id": thread_id, "workdir": str(resolved)}

    def get_working_directory(self, thread_id: str) -> str:
        return str(self._thread_workdirs.get(thread_id, DEFAULT_AGENT_WORKDIR))

    def clear_working_directory(self, thread_id: str) -> None:
        self._thread_workdirs.pop(thread_id, None)

    # ---------- dynamic configuration ----------
    def set_behavior_style(self, style: str):
        self.behavior_style = style or "standard"

    def set_extra_instruction(self, instruction: str):
        self.extra_instruction = instruction or ""

    def set_temperature(self, temperature: float):
        self.temperature = temperature
        try:
            llm.temperature = temperature
        except Exception:
            pass

    # ---------- graph nodes ----------
    def _planner_node(self, state: OrchestratorState):
        _log("WORKFLOW", "===== NODE: planner =====")
        triggering_content = state["messages"][-1].content
        goal = _as_text(triggering_content)
        attachments = _extract_image_blocks(triggering_content)
        _log(
            "WORKFLOW",
            f"🧭 Planning for goal: {goal!r}"
            + (f" (+{len(attachments)} image attachment(s))" if attachments else ""),
        )

        prior_history = state.get("conversation_history", []) or []
        history_text = _format_conversation_history(prior_history)
        history_note = (
            f"\n\nConversation so far (most recent {min(len(prior_history), MAX_HISTORY_TURNS_IN_PROMPT)} turns):\n{history_text}"
            if history_text else ""
        )

        style_note = f"\nDesired response style: {self.behavior_style}." if self.behavior_style else ""
        extra_note = f"\nAdditional instruction: {self.extra_instruction}" if self.extra_instruction else ""

        planning_prompt = f"""
You are a task planning system for a multi-agent orchestrator.
Break the goal below into a short sequence of atomic, actionable tasks.
For each task, assign an "agent_role": a short label for the kind of
specialist needed (e.g. "researcher", "calculator", "analyst", "writer").
Reuse the same role across tasks when it genuinely fits the same specialty.

If the conversation history below already contains the information needed to
answer the goal plan a
single task that answers directly from that history — do NOT plan tasks that
invent a way to "store" or "look up" data; you already have the transcript.
{history_note}

Goal: "{goal}"{style_note}{extra_note}

Return ONLY a JSON array, no prose, like:
[
  {{"id": "T1", "description": "...", "agent_role": "researcher"}},
  {{"id": "T2", "description": "...", "agent_role": "writer"}}
]

Use between 1 and {MAX_TASKS} tasks. Keep each task atomic.
"""
        _log_block("AI-REQUEST", "Prompt sent to LLM (planner)", planning_prompt)
        response = llm.invoke(
            planning_prompt,
            config={
                "run_name": "planner_llm",
                "tags": ["planner"],
                "metadata": {"goal": goal},
            },
        )
        _record_token_usage("planner", response, ["planner"])
        _log_block("AI-REPLY", "Raw LLM reply (planner)", _as_text(response.content))
        plan = parse_json_safely(response.content, default=None)
        if not plan or not isinstance(plan, list):
            plan = [{"id": "T1", "description": goal, "agent_role": "general"}]
        plan = plan[:MAX_TASKS]

        _log_block("WORKFLOW", "📋 Task plan", json.dumps(plan, indent=2))

        return {
            "goal": goal,
            "task_plan": plan,
            "current_task_idx": 0,
            "task_results": {},
            "task_messages": [],
            "retry_count": 0,
            "last_verdict": {},
            "conversation_history": [{"role": "user", "content": goal}],
            "attachments": attachments,
        }

    def _agent_executor_node(self, state: OrchestratorState):
        _log("WORKFLOW", "===== NODE: agent_executor =====")
        idx = state["current_task_idx"]
        task = state["task_plan"][idx]
        role = task["agent_role"]
        agent = self.agent_factory.get_or_create(role, task["description"])

        task_messages = state.get("task_messages") or []
        if not task_messages:
            context = self._build_context(state.get("task_results", {}), task)
            history_text = _format_conversation_history(state.get("conversation_history", []) or [])

            init_prompt = f"Task: {task['description']}"
            if history_text:
                init_prompt += f"\n\nFull conversation so far (use this as ground truth for anything the user has already said, e.g. their name or prior requests):\n{history_text}"
            if context:
                init_prompt += f"\n\nRelevant context from earlier tasks:\n{context}"

            # Regenerate the tool directive fresh so it reflects the agent's current tools
            full_system_prompt = agent["system_prompt"] + self.agent_factory.tool_directive(agent["tool_names"])
            attachments = state.get("attachments") or []
            if attachments:
                # Multimodal content: text block plus whatever images came in
                # with the triggering user message, so this task's agent can
                # actually see them, not just read a text description.
                init_content = [{"type": "text", "text": init_prompt}, *attachments]
            else:
                init_content = init_prompt
            task_messages = [
                SystemMessage(content=full_system_prompt),
                HumanMessage(content=init_content),
            ]

        tool_results_count = sum(1 for m in task_messages if isinstance(m, ToolMessage))
        if tool_results_count >= MAX_TOOL_CALLS_PER_TASK:
            _log("WORKFLOW", f"⛔ [{role}] task {task['id']}: tool call limit reached ({tool_results_count}/{MAX_TOOL_CALLS_PER_TASK}), forcing final answer")
            task_messages = task_messages + [
                HumanMessage(content="You have reached the maximum number of tool calls for this task. Do NOT invoke any more tools. Provide your final answer immediately based on the data retrieved so far.")
            ]

        _log("WORKFLOW", f"🤖 [{role}] executing task {task['id']}: {task['description']}")
        _log("AI-REQUEST", f"Invoking agent LLM [{role}] with {len(task_messages)} message(s) in context")
        response = agent["llm"].invoke(
            task_messages,
            config={
                "run_name": f"agent_exec:{role}",
                "tags": ["agent_executor", role, task["id"]],
                "metadata": {"role": role, "task_id": task["id"], "task_description": task["description"]},
            },
        )
        _record_token_usage(f"agent_executor:{role}:{task['id']}", response, ["agent_executor", role])

        reply_text = _as_text(response.content)
        _log_block("AI-REPLY", f"[{role}] LLM reply for task {task['id']}", reply_text or "(no text content)")
        tool_calls = getattr(response, "tool_calls", None) or []
        if tool_calls:
            for tc in tool_calls:
                _log("AI-REPLY", f"[{role}] requested tool call -> {tc['name']}({json.dumps(tc.get('args', {}), default=str)})")
        else:
            _log("AI-REPLY", f"[{role}] no tool calls requested; treating reply as this task's answer")

        task_messages = task_messages + [response]

        # Mirrored into top-level `messages` so LangSmith's Messages panel shows this turn
        return {"task_messages": task_messages, "messages": [response]}

    def _tools_node(self, state: OrchestratorState):
        """Custom tool executor operating on task_messages."""
        _log("WORKFLOW", "===== NODE: tools =====")
        idx = state["current_task_idx"]
        task = state["task_plan"][idx] if idx < len(state["task_plan"]) else {"id": "?", "agent_role": "?"}
        last = state["task_messages"][-1]
        tool_calls = getattr(last, "tool_calls", None) or []
        _log("WORKFLOW", f"🛠️ executing {len(tool_calls)} tool call(s) for task {task['id']}")
        tool_messages = []
        pending_images: list[dict] = []
        for call in tool_calls:
            tool_name = call["name"]
            tool_args = call.get("args", {})
            _log("TOOL-CALL", f"→ {tool_name}({json.dumps(tool_args, default=str)})")
            tool_obj = self.tool_registry.get_tool(tool_name)
            try:
                if tool_obj is None:
                    result = f"Error: tool '{tool_name}' not found"
                    _log("TOOL-RESULT", f"✗ {tool_name}: {result}")
                else:
                    result = tool_obj.invoke(
                        tool_args,
                        config={
                            "run_name": f"tool:{tool_name}",
                            "tags": ["tool_call", tool_name, task["id"]],
                            "metadata": {"task_id": task["id"], "agent_role": task.get("agent_role", "?")},
                        },
                    )
                    _log_block("TOOL-RESULT", f"✓ {tool_name} result", str(result))
            except Exception as e:
                result = f"Error executing tool '{tool_name}': {e}"
                _log("TOOL-RESULT", f"✗ {tool_name} raised: {e}")

            if isinstance(result, dict) and result.get("is_image"):
                # ToolMessage content must stay a plain string (OpenAI's API
                # doesn't accept image content there) -- so the ToolMessage
                # itself just confirms the load, and the actual image_url
                # block is queued to arrive as a follow-up HumanMessage the
                # agent sees on its very next turn.
                pending_images.append(
                    {"type": "image_url", "image_url": {"url": f"data:{result['mime_type']};base64,{result['data_b64']}"}}
                )
                summary = f"Loaded image '{result.get('path')}' ({result['mime_type']}); it is now attached for you to view."
                tool_messages.append(ToolMessage(content=summary, tool_call_id=call["id"]))
            else:
                tool_messages.append(ToolMessage(content=str(result), tool_call_id=call["id"]))

        updated_task_messages = state["task_messages"] + tool_messages
        extra_messages = list(tool_messages)
        if pending_images:
            image_message = HumanMessage(
                content=[{"type": "text", "text": "Here is the image you just loaded:"}, *pending_images]
            )
            updated_task_messages = updated_task_messages + [image_message]
            extra_messages = extra_messages + [image_message]

        # Mirrored into `messages` too, same reason as agent_executor_node above
        _log("WORKFLOW", f"🛠️ tool execution complete for task {task['id']}, returning to agent_executor")
        return {"task_messages": updated_task_messages, "messages": extra_messages}

    def _evaluator_node(self, state: OrchestratorState):
        _log("WORKFLOW", "===== NODE: evaluator =====")
        idx = state["current_task_idx"]
        task = state["task_plan"][idx]
        task_messages = state["task_messages"]
        final_content = _as_text(task_messages[-1].content) if task_messages else ""
        if not final_content and task_messages:
            for m in reversed(task_messages):
                if getattr(m, "content", None):
                    final_content = _as_text(m.content)
                    break

        eval_prompt = f"""
Task: {task['description']}
Agent output: {final_content}

Does this output satisfactorily complete the task? Respond with ONLY JSON:
{{"status": "PASS" or "RETRY", "reason": "short reason", "feedback": "what to fix if RETRY"}}
"""
        _log("AI-REQUEST", f"Invoking evaluator LLM for task {task['id']}")
        eval_response = llm.invoke(
            eval_prompt,
            config={
                "run_name": f"evaluator:{task['id']}",
                "tags": ["evaluator", task["agent_role"], task["id"]],
                "metadata": {"task_id": task["id"], "role": task["agent_role"]},
            },
        )
        _record_token_usage(f"evaluator:{task['id']}", eval_response, ["evaluator", task["agent_role"]])
        _log_block("AI-REPLY", f"Raw evaluator reply for task {task['id']}", _as_text(eval_response.content))
        verdict = parse_json_safely(
            eval_response.content,
            default={"status": "PASS", "reason": "auto-accepted (unparseable verdict)", "feedback": ""},
        )
        _log("WORKFLOW", f"🧪 Evaluation for {task['id']}: {verdict}")

        retry_count = state.get("retry_count", 0)
        if verdict.get("status") == "RETRY" and retry_count < MAX_RETRIES:
            _log("WORKFLOW", f"🔁 Task {task['id']} sent back for retry ({retry_count + 1}/{MAX_RETRIES}): {verdict.get('feedback', '')}")
            feedback_msg = HumanMessage(
                content=f"Evaluator feedback: {verdict.get('feedback', '')}. Please revise and try again."
            )
            return {
                "task_messages": state["task_messages"] + [feedback_msg],
                "messages": [feedback_msg],
                "retry_count": retry_count + 1,
                "last_verdict": verdict,
            }

        # Result accepted: either PASS, or retries exhausted
        _log("WORKFLOW", f"✅ Task {task['id']} accepted, moving to next task")
        results = dict(state.get("task_results", {}))
        results[task["id"]] = final_content
        summary_msg = AIMessage(
            content=f"[{task['agent_role']}] completed '{task['description']}':\n{final_content[:400]}"
        )
        return {
            "task_results": results,
            "current_task_idx": idx + 1,
            "task_messages": [],
            "retry_count": 0,
            "last_verdict": verdict,
            "messages": [summary_msg],
        }

    def _assembler_node(self, state: OrchestratorState):
        _log("WORKFLOW", "===== NODE: assembler =====")
        results = state.get("task_results", {})
        plan = state.get("task_plan", [])

        summary_prompt = f"Original goal: {state['goal']}\n\nTask results:\n"
        for t in plan:
            summary_prompt += f"- {t['description']}: {results.get(t['id'], 'N/A')}\n"
        summary_prompt += "\nWrite a single, coherent final answer to the original goal, using the results above."

        _log_block("AI-REQUEST", "Prompt sent to LLM (assembler)", summary_prompt)
        final = llm.invoke(
            summary_prompt,
            config={
                "run_name": "assembler_llm",
                "tags": ["assembler"],
                "metadata": {"goal": state["goal"]},
            },
        )
        _record_token_usage("assembler", final, ["assembler"])
        _log_block("FINAL-RESPONSE", "📦 Final response to user", _as_text(final.content), max_chars=10_000)
        return {
            "messages": [final],
            "conversation_history": [{"role": "assistant", "content": final.content}],
        }

    # ---------- routing ----------
    @staticmethod
    def _route_after_planner(state: OrchestratorState):
        route = "assembler" if state["current_task_idx"] >= len(state["task_plan"]) else "agent_executor"
        _log("WORKFLOW", f"↳ route_after_planner -> {route}")
        return route

    @staticmethod
    def _route_after_agent(state: OrchestratorState):
        task_messages = state.get("task_messages", [])
        if not task_messages:
            route = "evaluator"
        else:
            last = task_messages[-1]
            tool_call_count = sum(1 for m in task_messages if isinstance(m, ToolMessage))
            if getattr(last, "tool_calls", None) and tool_call_count < MAX_TOOL_CALLS_PER_TASK:
                route = "tools"
            else:
                route = "evaluator"
        _log("WORKFLOW", f"↳ route_after_agent -> {route}")
        return route

    @staticmethod
    def _route_after_evaluator(state: OrchestratorState):
        if state.get("task_messages"):  # retry pending, still has scratch messages
            route = "agent_executor"
        elif state["current_task_idx"] >= len(state["task_plan"]):
            route = "assembler"
        else:
            route = "agent_executor"
        _log("WORKFLOW", f"↳ route_after_evaluator -> {route}")
        return route

    # ---------- graph build ----------
    def _build_graph(self):
        graph = StateGraph(OrchestratorState)

        graph.add_node("planner", self._planner_node)
        graph.add_node("agent_executor", self._agent_executor_node)
        graph.add_node("tools", self._tools_node)
        graph.add_node("evaluator", self._evaluator_node)
        graph.add_node("assembler", self._assembler_node)

        graph.add_edge(START, "planner")
        graph.add_conditional_edges(
            "planner", self._route_after_planner,
            {"agent_executor": "agent_executor", "assembler": "assembler"},
        )
        graph.add_conditional_edges(
            "agent_executor", self._route_after_agent,
            {"tools": "tools", "evaluator": "evaluator"},
        )
        graph.add_edge("tools", "agent_executor")
        graph.add_conditional_edges(
            "evaluator", self._route_after_evaluator,
            {"agent_executor": "agent_executor", "assembler": "assembler"},
        )
        graph.add_edge("assembler", END)

        db_path = "DB/dynamic_chatbot.db"
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        conn = sqlite3.connect(database=db_path, check_same_thread=False)
        checkpointer = SqliteSaver(conn=conn)

        return graph.compile(checkpointer=checkpointer)

    # ---------- context sharing helper ----------
    @staticmethod
    def _build_context(task_results: dict, task: dict) -> str:
        """Pass only prior task results into a new task's prompt, keeping context small."""
        if not task_results:
            return ""
        lines = [f"- {tid}: {str(res)[:300]}" for tid, res in task_results.items()]
        return "\n".join(lines)

    # ---------- public API ----------
    def add_tool_from_prompt(self, prompt: str, tool_name: str):
        try:
            self.tool_registry.create_tool_from_prompt(prompt, tool_name)
            self.chatbot = self._build_graph()
            print(f"✅ Graph rebuilt with new tool: {tool_name}")
            return True
        except Exception as e:
            print(f"❌ Failed to add tool: {e}")
            return False

    def remove_tool(self, tool_name: str) -> bool:
        """
        Unregister a tool so no agent can call it and the LLM stops seeing it
        in "Available tools" (used by create_agent/refresh_tools when
        deciding whether to reuse or auto-create a tool). If a task later
        needs equivalent functionality again, the normal auto-tool-creation
        path (_create_missing_tools) will recreate it from scratch, since it
        no longer matches anything in list_tools().
        """
        removed = self.tool_registry.remove_tool(tool_name)
        if removed:
            self.agent_factory.remove_tool_from_agents(tool_name)
            self.chatbot = self._build_graph()
            print(f"✅ Graph rebuilt after removing tool: {tool_name}")
        return removed

    def get_tool_info(self) -> str:
        return json.dumps(self.tool_registry.list_tools(), indent=2)

    def get_agent_info(self) -> str:
        return json.dumps(self.agent_factory.list_agents(), indent=2)

    def get_artifacts(self) -> list[dict]:
        return self.tool_registry.list_artifacts()

    def list_threads(self, limit: int = 50) -> list[dict]:
        """List known thread_ids with a title (first user message) and last-updated time, most recent first."""
        latest_by_thread: dict[str, dict] = {}
        for cp_tuple in self.chatbot.checkpointer.list(None):
            thread_id = cp_tuple.config["configurable"]["thread_id"]
            ts = cp_tuple.checkpoint.get("ts", "")
            if thread_id in latest_by_thread and latest_by_thread[thread_id]["updated_at"] >= ts:
                continue
            messages = cp_tuple.checkpoint.get("channel_values", {}).get("messages", [])
            title = ""
            for m in messages:
                if getattr(m, "type", "") == "human" and getattr(m, "content", ""):
                    title = m.content
                    break
            latest_by_thread[thread_id] = {
                "thread_id": thread_id,
                "title": (title[:60] + "…") if len(title) > 60 else title,
                "updated_at": ts,
            }
        threads = sorted(latest_by_thread.values(), key=lambda t: t["updated_at"], reverse=True)
        return threads[:limit]

    def get_thread_history(self, thread_id: str) -> list[dict]:
        """Return a thread's full message history as plain role/content dicts."""
        state = self.chatbot.get_state({"configurable": {"thread_id": thread_id}})
        messages = state.values.get("messages", []) if state and state.values else []
        return [{"role": getattr(m, "type", "unknown"), "content": getattr(m, "content", "")} for m in messages]

    def run(self, user_input: str, thread_id: str, requirements: dict = None, doc_ids: list[str] = None,
            workdir: str = None):
        """
        requirements: {
            "new_tools": [{"name": "...", "prompt": "..."}],
            "dynamic_behavior": "concise" | "detailed" | ...,
            "preprocessing": "extra instruction text",
            "temperature": 0.0-1.0,
        }
        doc_ids: doc_ids returned by document_pipeline.process_upload() for
            any files (documents or images) the user attached to this
            message. Text-extractable docs are folded into the prompt;
            images are attached as real vision content -- see
            document_pipeline.build_multimodal_message.
        workdir: an on-disk folder to point this thread's local file tools
            (read_file/write_file/list_directory/view_image/create_artifact)
            at, same as select_working_directory(). Accepted here too so a
            caller can select a local working directory in the SAME request
            that starts a brand-new thread, instead of needing an already-
            known thread_id to call select_working_directory() with first.
            Raises WorkdirSelectionError if the path doesn't exist, isn't a
            directory, or is refused (see set_working_directory).
        """
        if workdir:
            self.set_working_directory(thread_id, workdir)

        if requirements:
            if requirements.get("new_tools"):
                for tool_spec in requirements["new_tools"]:
                    self.add_tool_from_prompt(tool_spec["prompt"], tool_spec["name"])
            if "dynamic_behavior" in requirements:
                self.set_behavior_style(requirements["dynamic_behavior"])
            if "preprocessing" in requirements:
                self.set_extra_instruction(requirements["preprocessing"])
            if "temperature" in requirements:
                self.set_temperature(requirements["temperature"])

        config = {
            "configurable": {"thread_id": thread_id},
            "run_name": "orchestrator_run",
            "tags": ["orchestrator_run"],
            "metadata": {"goal": user_input, "thread_id": thread_id},
        }

        message_content = docpipe.build_multimodal_message(user_input, doc_ids)

        _log("WORKFLOW", f"########## RUN START (thread={thread_id}) ##########")
        _log("WORKFLOW", f"User input: {user_input!r}")

        usage_records: list = []
        ctx_token = _token_usage_ctx.set(usage_records)
        workdir_token = _workdir_ctx.set(self._thread_workdirs.get(thread_id))
        try:
            response = self.chatbot.invoke(
                {"messages": [HumanMessage(content=message_content)]},
                config=config,
            )
        finally:
            _token_usage_ctx.reset(ctx_token)
            _workdir_ctx.reset(workdir_token)

        token_usage = self._summarize_token_usage(usage_records)
        self._print_token_usage(user_input, token_usage)
        response["token_usage"] = token_usage

        final_ai_text = ""
        for m in reversed(response.get("messages", [])):
            if isinstance(m, AIMessage) and getattr(m, "content", None):
                final_ai_text = _as_text(m.content)
                break
        _log_block("FINAL-RESPONSE", f"########## RUN END (thread={thread_id}) — response returned to user", final_ai_text, max_chars=10_000)
        return response

    # ---------- token usage helpers ----------
    @staticmethod
    def _summarize_token_usage(records: list[dict]) -> dict:
        """Turn the raw list of per-LLM-call token records collected during one run() into a
        step-by-step breakdown plus totals, e.g.:
        {
          "calls": [
            {"step": 1, "node": "planner", "tags": [...], "input_tokens": 320, "output_tokens": 85, "total_tokens": 405},
            ...
          ],
          "totals": {"input_tokens": ..., "output_tokens": ..., "total_tokens": ..., "llm_calls": N}
        }
        """
        calls = []
        totals = {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
        for i, r in enumerate(records, start=1):
            calls.append({
                "step": i,
                "node": r["node"],
                "tags": r.get("tags", []),
                "input_tokens": r["input_tokens"],
                "output_tokens": r["output_tokens"],
                "total_tokens": r["total_tokens"],
            })
            totals["input_tokens"] += r["input_tokens"]
            totals["output_tokens"] += r["output_tokens"]
            totals["total_tokens"] += r["total_tokens"]
        totals["llm_calls"] = len(records)
        return {"calls": calls, "totals": totals}

    @staticmethod
    def _print_token_usage(user_input: str, token_usage: dict):
        totals = token_usage["totals"]
        print(f"\n🔢 Token usage for prompt: {user_input[:80]!r}")
        for c in token_usage["calls"]:
            print(
                f"   [{c['step']}] {c['node']:<32} "
                f"in={c['input_tokens']:<6} out={c['output_tokens']:<6} total={c['total_tokens']}"
            )
        print(
            f"   {'TOTAL (' + str(totals['llm_calls']) + ' calls)':<36} "
            f"in={totals['input_tokens']:<6} out={totals['output_tokens']:<6} total={totals['total_tokens']}\n"
        )




agent_manager = DynamicAgentManager()


def get_agent_tools() -> str:
    return agent_manager.get_tool_info()


def get_agent_registry() -> str:
    return agent_manager.get_agent_info()


def get_artifacts() -> list[dict]:
    return agent_manager.get_artifacts()


def add_tool_dynamically(tool_name: str, tool_prompt: str) -> bool:
    return agent_manager.add_tool_from_prompt(tool_prompt, tool_name)


def remove_tool_dynamically(tool_name: str) -> bool:
    return agent_manager.remove_tool(tool_name)


def run_agent_with_requirements(user_input: str, thread_id: str, requirements: dict = None, doc_ids: list[str] = None,
                                 workdir: str = None):
    return agent_manager.run(user_input, thread_id, requirements, doc_ids, workdir)


def select_working_directory(thread_id: str, path: str) -> dict:
    """Cowork-style folder picker: point a thread's file tools at an
    existing directory on disk. Raises WorkdirSelectionError on an
    invalid or refused path."""
    return agent_manager.set_working_directory(thread_id, path)


def get_working_directory(thread_id: str) -> str:
    return agent_manager.get_working_directory(thread_id)


def clear_working_directory(thread_id: str) -> None:
    """Reset a thread's file tools back to DEFAULT_AGENT_WORKDIR."""
    agent_manager.clear_working_directory(thread_id)


def upload_file(filename: str, data: bytes) -> str:
    """Process an uploaded file (document or image) and return its doc_id,
    ready to pass into run_agent_with_requirements(..., doc_ids=[...])."""
    return docpipe.process_upload(filename, data)